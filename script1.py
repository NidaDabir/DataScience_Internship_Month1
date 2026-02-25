from pptx import Presentation
from pptx.util import Inches, Pt


def create_presentation():
    prs = Presentation()

    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Object Counter using OpenCV"
    subtitle.text = "Name: Dabir Nida Khalid\nMonth 2 Task 1\n\nRepo: https://github.com/NidaDabir/DataScience_Internship_Month1.git"

    # Data for the rest of the slides
    slides_data = [
        ("Objective", [
            "To detect and count objects in video",
            "Using OpenCV and Python"
        ]),
        ("Dataset", [
            "Used sample video",
            "Contains moving people"
        ]),
        ("Steps Performed", [
            "Loaded video using OpenCV",
            "Converted to grayscale",
            "Applied threshold",
            "Created mask",
            "Counted objects",
            "Displayed result"
        ]),
        ("Output", [
            "Object count displayed on screen",
            "Video runs successfully",
            "(Note: Please insert your output screenshot here)"
        ]),
        ("Tools Used", [
            "Python",
            "OpenCV",
            "PyCharm"
        ]),
        ("Conclusion", [
            "Successfully counted objects",
            "Learned OpenCV basics"
        ])
    ]

    # --- Loop to create Slides 2-7 ---
    for slide_title, points in slides_data:
        slide = prs.slides.add_slide(bullet_slide_layout)
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]

        title_shape.text = slide_title
        tf = body_shape.text_frame

        for point in points:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(20)  # Adjust font size if needed

    # Save the presentation
    prs.save('Month2_Task1_Object_Counter.pptx')
    print("Presentation saved successfully as 'Month2_Task1_Object_Counter.pptx'")


if __name__ == "__main__":
    create_presentation()